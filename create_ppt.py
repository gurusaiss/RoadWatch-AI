"""
RoadWatch PPT Generator — IIT Madras National Road Safety Hackathon 2026
Generates a professional 7-slide PowerPoint presentation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
import os

# ── Palette ────────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
BLUE        = RGBColor(0x1A, 0x6B, 0xB5)   # brand blue
CYAN        = RGBColor(0x00, 0xC9, 0xFF)   # accent cyan
ORANGE      = RGBColor(0xFF, 0x6B, 0x35)   # accent orange
GREEN       = RGBColor(0x2E, 0xCC, 0x71)   # success green
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xF0, 0xF4, 0xF8)
MID_GRAY    = RGBColor(0x8A, 0x9A, 0xAD)
DARK_TEXT   = RGBColor(0x1A, 0x1A, 0x2E)
YELLOW      = RGBColor(0xFF, 0xD7, 0x00)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── Helpers ─────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]   # truly blank
    return prs.slides.add_slide(blank_layout)

def fill_bg(slide, color: RGBColor):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox

def add_para(tf, text, font_size=14, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, space_before=6):
    from pptx.util import Pt as PT
    from pptx.oxml.ns import qn
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = PT(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return p

def accent_bar(slide, color=CYAN, height=Inches(0.06)):
    """Thin accent bar at top of slide."""
    add_rect(slide, 0, 0, SLIDE_W, height, color)

def slide_number(slide, num, total=7):
    add_text(slide, f"{num} / {total}",
             SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.45),
             Inches(1.0), Inches(0.35),
             font_size=9, color=MID_GRAY, align=PP_ALIGN.RIGHT)

def add_card(slide, left, top, width, height, bg=RGBColor(0x15,0x2A,0x3E),
             radius=False):
    r = add_rect(slide, left, top, width, height, bg)
    return r

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title / Welcome
# ═══════════════════════════════════════════════════════════════════════════
def slide1_title(prs):
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)

    # gradient-like layered rectangles
    add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    add_rect(sl, 0, 0, Inches(6), SLIDE_H, RGBColor(0x0A, 0x14, 0x20))

    # bold left accent bar
    add_rect(sl, 0, 0, Inches(0.18), SLIDE_H, CYAN)

    # top-right decorative circle (fake with wide rect, simulated)
    add_rect(sl, Inches(9.5), Inches(-1), Inches(5), Inches(5),
             RGBColor(0x1A, 0x3A, 0x5C))

    # RoadWatch logo-text
    add_text(sl, "RoadWatch",
             Inches(0.5), Inches(1.2), Inches(7), Inches(1.6),
             font_size=72, bold=True, color=WHITE)

    add_text(sl, "AI-Powered Road Safety & Transparency Platform",
             Inches(0.5), Inches(2.75), Inches(8.5), Inches(0.9),
             font_size=24, bold=False, color=CYAN)

    add_text(sl, "IIT Madras — National Road Safety Hackathon 2026",
             Inches(0.5), Inches(3.65), Inches(8), Inches(0.5),
             font_size=15, color=LIGHT_GRAY, italic=True)

    # divider line
    add_rect(sl, Inches(0.5), Inches(4.25), Inches(5.5), Inches(0.04), CYAN)

    # team info
    add_text(sl, "Team  |  RoadWatch Track",
             Inches(0.5), Inches(4.4), Inches(6), Inches(0.4),
             font_size=14, color=MID_GRAY)

    # right panel — 3 quick stats
    for i, (val, lbl) in enumerate([
        ("35",  "Roads Monitored"),
        ("13",  "Countries Covered"),
        ("AI",  "Damage Assessment"),
    ]):
        y = Inches(1.8 + i * 1.6)
        add_rect(sl, Inches(10.0), y, Inches(2.8), Inches(1.2),
                 RGBColor(0x15, 0x2A, 0x3E))
        add_text(sl, val,
                 Inches(10.1), y + Inches(0.08), Inches(2.6), Inches(0.6),
                 font_size=34, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
        add_text(sl, lbl,
                 Inches(10.1), y + Inches(0.62), Inches(2.6), Inches(0.42),
                 font_size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    accent_bar(sl, CYAN)
    slide_number(sl, 1)
    return sl


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem & National Impact
# ═══════════════════════════════════════════════════════════════════════════
def slide2_problem(prs):
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)
    accent_bar(sl, ORANGE)

    # header
    add_rect(sl, 0, Inches(0.06), SLIDE_W, Inches(1.15), RGBColor(0x0D,0x22,0x38))
    add_text(sl, "The Problem  —  India's Road Safety Crisis",
             Inches(0.5), Inches(0.18), Inches(11), Inches(0.8),
             font_size=28, bold=True, color=WHITE)

    # 4 stat cards
    stats = [
        ("1,68,491", "Road accident deaths\nin India (2022)", ORANGE),
        ("₹5.96L Cr", "Annual road network\nmaintenance budget", BLUE),
        ("~40%", "Complaints never\nreached right authority", RGBColor(0xE7,0x4C,0x3C)),
        ("63%", "Roads lack\ntransparent fund tracking", RGBColor(0xFF,0xA5,0x00)),
    ]
    for i, (val, lbl, col) in enumerate(stats):
        x = Inches(0.35 + i * 3.2)
        add_rect(sl, x, Inches(1.45), Inches(3.0), Inches(1.8),
                 RGBColor(0x10, 0x24, 0x36))
        add_rect(sl, x, Inches(1.45), Inches(3.0), Inches(0.07), col)
        add_text(sl, val, x + Inches(0.12), Inches(1.6),
                 Inches(2.76), Inches(0.75),
                 font_size=30, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_text(sl, lbl, x + Inches(0.08), Inches(2.3),
                 Inches(2.84), Inches(0.72),
                 font_size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # pain-point bullets
    add_text(sl, "Root Causes",
             Inches(0.5), Inches(3.45), Inches(4), Inches(0.45),
             font_size=17, bold=True, color=CYAN)
    pains = [
        "No single platform shows road contractor + EE + budget in one place",
        "Citizens file complaints with wrong departments — months of delay",
        "Budget anomalies (under-utilisation) invisible to public & auditors",
        "No offline access in rural / low-connectivity areas",
        "Existing portals are India-only with no global benchmarking",
    ]
    for j, p in enumerate(pains):
        add_rect(sl, Inches(0.5), Inches(3.98 + j * 0.58),
                 Inches(0.18), Inches(0.38), ORANGE)
        add_text(sl, p,
                 Inches(0.82), Inches(3.94 + j * 0.58),
                 Inches(5.6), Inches(0.46),
                 font_size=12.5, color=LIGHT_GRAY)

    # right image-like panel
    add_rect(sl, Inches(6.8), Inches(3.35), Inches(6.2), Inches(3.85),
             RGBColor(0x10, 0x24, 0x36))
    add_text(sl, "Status Quo",
             Inches(7.0), Inches(3.5), Inches(5.8), Inches(0.5),
             font_size=16, bold=True, color=MID_GRAY, align=PP_ALIGN.CENTER)
    issues = [
        ("iRAD / NIC portals", "No budget transparency"),
        ("CPGRAMS", "Generic — not road-specific"),
        ("State PWD portals", "Siloed, no cross-state view"),
        ("Manual complaint filing", "Wrong routing, no tracking"),
    ]
    for k, (sys, prob) in enumerate(issues):
        y = Inches(4.1 + k * 0.7)
        add_rect(sl, Inches(7.0), y, Inches(2.6), Inches(0.52),
                 RGBColor(0x1A, 0x32, 0x4A))
        add_rect(sl, Inches(9.8), y, Inches(3.0), Inches(0.52),
                 RGBColor(0x2A, 0x12, 0x12))
        add_text(sl, sys, Inches(7.1), y + Inches(0.07), Inches(2.4), Inches(0.38),
                 font_size=11, bold=True, color=WHITE)
        add_text(sl, prob, Inches(9.9), y + Inches(0.07), Inches(2.8), Inches(0.38),
                 font_size=11, color=RGBColor(0xFF,0x8A,0x8A))

    slide_number(sl, 2)
    return sl


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Gap Analysis & Our Differentiation
# ═══════════════════════════════════════════════════════════════════════════
def slide3_gap(prs):
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)
    accent_bar(sl, YELLOW)

    add_rect(sl, 0, Inches(0.06), SLIDE_W, Inches(1.15), RGBColor(0x0D,0x22,0x38))
    add_text(sl, "Gap in Existing Systems  →  RoadWatch Advantage",
             Inches(0.5), Inches(0.18), Inches(12), Inches(0.8),
             font_size=28, bold=True, color=WHITE)

    # comparison table header
    headers = ["Feature", "iRAD / NIC", "CPGRAMS", "State PWDs", "RoadWatch ✓"]
    col_w   = [Inches(2.8), Inches(2.1), Inches(2.1), Inches(2.1), Inches(2.3)]
    col_x   = [Inches(0.3)]
    for w in col_w[:-1]:
        col_x.append(col_x[-1] + w + Inches(0.08))

    # header row
    for i, (h, x, w) in enumerate(zip(headers, col_x, col_w)):
        bg = BLUE if i == 0 else (GREEN if i == 4 else RGBColor(0x1A,0x32,0x4A))
        add_rect(sl, x, Inches(1.4), w, Inches(0.5), bg)
        add_text(sl, h, x + Inches(0.08), Inches(1.45), w - Inches(0.1), Inches(0.38),
                 font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    rows = [
        ("Road contractor + EE info", "✗", "✗", "Partial", "✔ Full detail"),
        ("Budget sanctioned vs spent", "✗", "✗", "✗", "✔ With source"),
        ("Smart complaint routing", "Partial", "Generic", "✗", "✔ AI-routed"),
        ("Multi-country support",    "✗", "✗", "✗", "✔ 13 countries"),
        ("AI damage assessment",     "✗", "✗", "✗", "✔ CV + LLM"),
        ("Offline PWA capability",   "✗", "✗", "✗", "✔ Service Worker"),
        ("Budget anomaly alerts",    "✗", "✗", "✗", "✔ Auto-flagged"),
    ]
    for r, row in enumerate(rows):
        y = Inches(1.98 + r * 0.67)
        bg_row = RGBColor(0x10,0x22,0x34) if r % 2 == 0 else RGBColor(0x0D,0x1C,0x2C)
        for i, (cell, x, w) in enumerate(zip(row, col_x, col_w)):
            bg = bg_row
            fc = LIGHT_GRAY
            if i == 0:
                bg = RGBColor(0x15, 0x2A, 0x40)
                fc = WHITE
            elif i == 4:
                bg = RGBColor(0x0D, 0x2B, 0x1A)
                fc = GREEN
            elif cell == "✗":
                fc = RGBColor(0xFF, 0x6B, 0x6B)
            add_rect(sl, x, y, w, Inches(0.58), bg)
            add_text(sl, cell, x + Inches(0.08), y + Inches(0.08),
                     w - Inches(0.1), Inches(0.42),
                     font_size=11.5, color=fc, align=PP_ALIGN.CENTER,
                     bold=(i == 4))

    slide_number(sl, 3)
    return sl


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Our Solution: RoadWatch Platform
# ═══════════════════════════════════════════════════════════════════════════
def slide4_solution(prs):
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)
    accent_bar(sl, GREEN)

    add_rect(sl, 0, Inches(0.06), SLIDE_W, Inches(1.15), RGBColor(0x0D,0x22,0x38))
    add_text(sl, "RoadWatch  —  One Platform. Complete Transparency.",
             Inches(0.5), Inches(0.18), Inches(12), Inches(0.8),
             font_size=28, bold=True, color=WHITE)

    # 6 feature cards in 2 rows of 3
    features = [
        ("🗺️ Interactive Map", BLUE,
         "Leaflet.js + MarkerCluster\nColour-coded condition overlays\nGeolocation & radius search"),
        ("📊 Budget Analytics", RGBColor(0x8E,0x44,0xAD),
         "Sanctioned vs spent tracking\nAnomaly auto-flag (<65%)\nYear-on-year charts"),
        ("🤖 AI Chatbot", CYAN,
         "Groq LLaMA 3.3 70B backend\nRule-based + LLM fallback\nMarkdown-rich responses"),
        ("📋 Smart Complaint Wizard", ORANGE,
         "4-step guided filing\nAI routes to correct EE\nReal-time status tracking"),
        ("📸 Damage Assessment", RGBColor(0xE7,0x4C,0x3C),
         "Image upload → CV analysis\nSeverity & damage-type output\nAuto-attached to complaint"),
        ("🌐 Global Coverage", GREEN,
         "13 countries: IN, UK, US, BD, ZA, MG, PL, DE, AU, NG, KE, JP, MY\nCurrency-aware display\nLocalised authority routing"),
    ]

    for i, (title, col, desc) in enumerate(features):
        row, col_i = divmod(i, 3)
        x = Inches(0.35 + col_i * 4.3)
        y = Inches(1.5 + row * 2.85)
        add_rect(sl, x, y, Inches(4.1), Inches(2.6), RGBColor(0x10,0x24,0x36))
        add_rect(sl, x, y, Inches(4.1), Inches(0.08), col)
        add_text(sl, title, x + Inches(0.15), y + Inches(0.15),
                 Inches(3.8), Inches(0.55),
                 font_size=15, bold=True, color=WHITE)
        add_text(sl, desc, x + Inches(0.15), y + Inches(0.72),
                 Inches(3.8), Inches(1.72),
                 font_size=11.5, color=LIGHT_GRAY)

    slide_number(sl, 4)
    return sl


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Architecture & Workflow
# ═══════════════════════════════════════════════════════════════════════════
def slide5_architecture(prs):
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)
    accent_bar(sl, BLUE)

    add_rect(sl, 0, Inches(0.06), SLIDE_W, Inches(1.15), RGBColor(0x0D,0x22,0x38))
    add_text(sl, "Architecture & Complaint Workflow",
             Inches(0.5), Inches(0.18), Inches(12), Inches(0.8),
             font_size=28, bold=True, color=WHITE)

    # ── Architecture diagram (left half) ──
    add_text(sl, "Tech Stack",
             Inches(0.4), Inches(1.38), Inches(5), Inches(0.45),
             font_size=15, bold=True, color=CYAN)

    layers = [
        ("FRONTEND (Browser / PWA)",   BLUE,
         "Vanilla JS · Leaflet.js · Chart.js · marked.js\nService Worker · Offline cache · Responsive CSS"),
        ("BACKEND API  (FastAPI / Python)", RGBColor(0x8E,0x44,0xAD),
         "REST endpoints · SQLAlchemy ORM · Pydantic v2\nHaversine geo-search · File uploads · aiofiles"),
        ("AI LAYER",  CYAN,
         "Groq LLaMA-3.3-70B · OpenAI GPT-4 Vision\nRule-based fallback · Damage CV assessment"),
        ("DATA STORE  (SQLite)",  RGBColor(0x27,0xAE,0x60),
         "Roads · Authorities · Complaints · Budget history\n35 roads · 37 authorities · 13 countries"),
    ]

    for i, (lbl, col, desc) in enumerate(layers):
        y = Inches(1.92 + i * 1.3)
        add_rect(sl, Inches(0.35), y, Inches(6.0), Inches(1.15),
                 RGBColor(0x10, 0x24, 0x36))
        add_rect(sl, Inches(0.35), y, Inches(0.18), Inches(1.15), col)
        add_text(sl, lbl, Inches(0.65), y + Inches(0.08),
                 Inches(5.5), Inches(0.42),
                 font_size=12, bold=True, color=col)
        add_text(sl, desc, Inches(0.65), y + Inches(0.5),
                 Inches(5.5), Inches(0.6),
                 font_size=10.5, color=LIGHT_GRAY)

    # connector arrows between layers
    for i in range(3):
        y_arrow = Inches(3.05 + i * 1.3)
        add_rect(sl, Inches(3.1), y_arrow, Inches(0.18), Inches(0.12),
                 RGBColor(0x4A, 0x6A, 0x8A))

    # ── Complaint flow (right half) ──
    add_text(sl, "Smart Complaint Flow",
             Inches(7.0), Inches(1.38), Inches(6), Inches(0.45),
             font_size=15, bold=True, color=ORANGE)

    steps = [
        ("1", "Citizen files complaint\n(GPS + image + description)", ORANGE),
        ("2", "AI CV assesses damage\n(type, severity, confidence)", CYAN),
        ("3", "AI routes to correct\nExecutive Engineer + dept.", BLUE),
        ("4", "EE receives notification;\nstatus tracked in real-time", GREEN),
        ("5", "Resolution logged; citizen\nnotified; data feeds analytics", RGBColor(0x9B,0x59,0xB6)),
    ]

    for i, (num, txt, col) in enumerate(steps):
        y = Inches(1.9 + i * 1.05)
        # circle number
        add_rect(sl, Inches(7.1), y, Inches(0.55), Inches(0.55), col)
        add_text(sl, num, Inches(7.1), y, Inches(0.55), Inches(0.55),
                 font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(sl, Inches(7.75), y, Inches(5.2), Inches(0.82),
                 RGBColor(0x10, 0x24, 0x36))
        add_text(sl, txt, Inches(7.9), y + Inches(0.06),
                 Inches(4.9), Inches(0.7),
                 font_size=11.5, color=LIGHT_GRAY)
        # connector line (except last)
        if i < 4:
            add_rect(sl, Inches(7.33), y + Inches(0.55), Inches(0.08),
                     Inches(0.5), RGBColor(0x4A,0x6A,0x8A))

    slide_number(sl, 5)
    return sl


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Impact, Scalability & Why RoadWatch Wins
# ═══════════════════════════════════════════════════════════════════════════
def slide6_impact(prs):
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)
    accent_bar(sl, RGBColor(0x9B,0x59,0xB6))

    add_rect(sl, 0, Inches(0.06), SLIDE_W, Inches(1.15), RGBColor(0x0D,0x22,0x38))
    add_text(sl, "Impact, Scalability & Why RoadWatch Wins",
             Inches(0.5), Inches(0.18), Inches(12), Inches(0.8),
             font_size=28, bold=True, color=WHITE)

    # ── Evaluation criteria coverage (left) ──
    add_text(sl, "Evaluation Criteria  (Rulebook)",
             Inches(0.4), Inches(1.38), Inches(6.5), Inches(0.45),
             font_size=15, bold=True, color=CYAN)

    criteria = [
        ("Data Accuracy",
         "Verified contractor, EE, relaying date, budget — all with source labels"),
        ("Correct Complaint Routing",
         "AI + DB authority lookup → routes to right EE/department always"),
        ("Budget Transparency",
         "Sanctioned vs spent per year with anomaly auto-flag & source URL"),
        ("Global Applicability",
         "13 countries, localised authorities, multi-currency budget display"),
        ("Offline Functionality",
         "Service Worker PWA: roads + stats cached; works without internet"),
        ("UI & Accessibility",
         "Responsive, hamburger-mobile, WCAG contrast, screen-reader labels"),
    ]
    for i, (crit, detail) in enumerate(criteria):
        y = Inches(1.92 + i * 0.86)
        add_rect(sl, Inches(0.35), y, Inches(6.2), Inches(0.76),
                 RGBColor(0x0D, 0x22, 0x38))
        add_rect(sl, Inches(0.35), y, Inches(0.12), Inches(0.76), GREEN)
        add_text(sl, f"✔  {crit}", Inches(0.55), y + Inches(0.04),
                 Inches(5.8), Inches(0.32),
                 font_size=12, bold=True, color=GREEN)
        add_text(sl, detail, Inches(0.55), y + Inches(0.36),
                 Inches(5.8), Inches(0.36),
                 font_size=10.5, color=LIGHT_GRAY)

    # ── Impact numbers (right top) ──
    add_text(sl, "Projected Impact",
             Inches(7.0), Inches(1.38), Inches(6), Inches(0.45),
             font_size=15, bold=True, color=ORANGE)

    impacts = [
        ("60%", "Faster complaint resolution\nwith AI routing"),
        ("100%", "Budget data\ntransparency"),
        ("13", "Countries on\nDay 1"),
        ("0", "Wrong-department\nmisdirections"),
    ]
    for i, (val, lbl) in enumerate(impacts):
        col_i = i % 2
        row_i = i // 2
        x = Inches(7.1 + col_i * 3.0)
        y = Inches(1.95 + row_i * 1.6)
        add_rect(sl, x, y, Inches(2.7), Inches(1.35),
                 RGBColor(0x10, 0x24, 0x36))
        add_text(sl, val, x, y + Inches(0.08), Inches(2.7), Inches(0.7),
                 font_size=36, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
        add_text(sl, lbl, x, y + Inches(0.72), Inches(2.7), Inches(0.55),
                 font_size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # ── Scalability note ──
    add_rect(sl, Inches(6.8), Inches(5.4), Inches(6.2), Inches(1.75),
             RGBColor(0x0D, 0x22, 0x38))
    add_text(sl, "Scale-up Roadmap",
             Inches(7.0), Inches(5.5), Inches(5.8), Inches(0.4),
             font_size=13, bold=True, color=CYAN)
    add_text(sl,
             "• Replace SQLite with PostgreSQL for production scale\n"
             "• Integrate PMGSY / NRIDA real data APIs\n"
             "• Mobile app (React Native) with full offline sync\n"
             "• Bhashini integration for 22 Indian language support",
             Inches(7.0), Inches(5.95), Inches(5.8), Inches(1.1),
             font_size=10.5, color=LIGHT_GRAY)

    slide_number(sl, 6)
    return sl


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Thank You
# ═══════════════════════════════════════════════════════════════════════════
def slide7_thankyou(prs):
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)

    add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    add_rect(sl, 0, 0, Inches(0.18), SLIDE_H, GREEN)
    add_rect(sl, Inches(7.5), 0, Inches(5.9), SLIDE_H, RGBColor(0x0A,0x18,0x26))

    # large thank you
    add_text(sl, "Thank You",
             Inches(0.5), Inches(1.5), Inches(7), Inches(1.6),
             font_size=72, bold=True, color=WHITE)

    add_text(sl, "Making India's roads safer — one data point at a time.",
             Inches(0.5), Inches(3.1), Inches(7), Inches(0.75),
             font_size=20, color=CYAN, italic=True)

    add_rect(sl, Inches(0.5), Inches(3.95), Inches(5.5), Inches(0.05), GREEN)

    add_text(sl, "RoadWatch  |  IIT Madras National Road Safety Hackathon 2026",
             Inches(0.5), Inches(4.15), Inches(7), Inches(0.45),
             font_size=14, color=MID_GRAY)

    # QR placeholder / demo info
    add_rect(sl, Inches(8.5), Inches(2.0), Inches(4.0), Inches(3.8),
             RGBColor(0x10, 0x24, 0x36))
    add_text(sl, "Platform Highlights",
             Inches(8.7), Inches(2.15), Inches(3.6), Inches(0.5),
             font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    highlights = [
        "🗺️  Live interactive road map",
        "📊  Budget anomaly detection",
        "🤖  AI chatbot (LLaMA / GPT-4)",
        "📸  Computer Vision damage CV",
        "🌐  13-country global support",
        "📱  Offline PWA capability",
    ]
    for i, h in enumerate(highlights):
        add_rect(sl, Inches(8.65), Inches(2.75 + i * 0.48),
                 Inches(3.65), Inches(0.4),
                 RGBColor(0x15, 0x2A, 0x3E))
        add_text(sl, h, Inches(8.8), Inches(2.78 + i * 0.48),
                 Inches(3.4), Inches(0.36),
                 font_size=11.5, color=LIGHT_GRAY)

    accent_bar(sl, GREEN)
    slide_number(sl, 7)
    return sl


# ═══════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════
if __name__ == "__main__":
    prs = new_prs()

    slide1_title(prs)
    slide2_problem(prs)
    slide3_gap(prs)
    slide4_solution(prs)
    slide5_architecture(prs)
    slide6_impact(prs)
    slide7_thankyou(prs)

    out_dir = r"D:\C\CODING\Hackathon\Road Safety\roadwatch"
    out_path = os.path.join(out_dir, "RoadWatch_Presentation.pptx")
    prs.save(out_path)
    print(f"Saved: {out_path}")
    print(f"Slides: {len(prs.slides)}")
